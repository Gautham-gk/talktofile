"""
Slide Agent — generates a PowerPoint presentation from document content.
Uses python-pptx to produce a .pptx file.
"""

import io
import json
from openai import AsyncOpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from core.session_store import DocumentData
from core.config import get_settings


_SYSTEM = """You are a presentation designer. Given document content, create a structured slide deck.

Generate 8-12 slides covering the key content. Each slide should have:
- A clear, concise title (max 8 words)
- 3-5 bullet points (each max 15 words)
- Optional speaker_note (1-2 sentences for the presenter)

The first slide is always a title slide with just a title and subtitle.
The last slide is always a summary/takeaways slide.

Return ONLY a JSON array (no markdown, no extra text):
[
  {"type": "title", "title": "...", "subtitle": "...", "speaker_note": "..."},
  {"type": "content", "title": "...", "bullets": ["...", "..."], "speaker_note": "..."},
  ...
]"""


def _build_context(documents: list[DocumentData]) -> str:
    parts = []
    for doc in documents:
        if isinstance(doc.summary, dict):
            s = doc.summary
            parts.append(
                f"=== {doc.filename} ===\n"
                f"Overview: {s.get('overview','')}\n"
                f"Key points: {'; '.join(s.get('key_points', []))}\n"
                f"Topics: {', '.join(s.get('topics', []))}"
            )
        if doc.chunks:
            content = "\n\n".join(doc.chunks)[:8000]
            parts.append(f"Content:\n{content}")
    return "\n\n".join(parts)


async def generate_slides_data(documents: list[DocumentData]) -> list[dict]:
    settings = get_settings()
    client = AsyncOpenAI(api_key=settings.openai_api_key)
    context = _build_context(documents)

    response = await client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": _SYSTEM},
            {"role": "user", "content": f"Document content:\n\n{context}"},
        ],
        temperature=0.4,
        max_tokens=3000,
    )

    raw = (response.choices[0].message.content or "[]").strip()
    if raw.startswith("```"):
        raw = raw.split("\n", 1)[-1].rsplit("```", 1)[0].strip()

    try:
        slides = json.loads(raw)
        return slides if isinstance(slides, list) else []
    except json.JSONDecodeError:
        return []


_BRAND_RED = RGBColor(0xE6, 0x00, 0x26)
_DARK = RGBColor(0x30, 0x30, 0x30)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT_GRAY = RGBColor(0xF8, 0xFA, 0xFC)


def _set_text(tf, text: str, size_pt: int, bold: bool = False, color: RGBColor = None):
    tf.text = text
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            if color:
                run.font.color.rgb = color


def build_pptx(slides_data: list[dict], doc_title: str = "Document") -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # completely blank

    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        slide_type = slide_data.get("type", "content")

        if slide_type == "title":
            # Red background for title slide
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = _BRAND_RED

            title_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.33), Inches(1.5))
            tf = title_box.text_frame
            tf.word_wrap = True
            _set_text(tf, slide_data.get("title", doc_title), 44, bold=True, color=_WHITE)

            subtitle = slide_data.get("subtitle", "")
            if subtitle:
                sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.9), Inches(11.33), Inches(1))
                tf2 = sub_box.text_frame
                tf2.word_wrap = True
                _set_text(tf2, subtitle, 24, bold=False, color=RGBColor(0xFF, 0xCC, 0xCC))
        else:
            # White background
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = _WHITE

            # Red accent bar at top
            accent = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                Inches(0), Inches(0), Inches(13.33), Inches(0.08)
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = _BRAND_RED
            accent.line.fill.background()

            # Title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.33), Inches(0.9))
            tf = title_box.text_frame
            tf.word_wrap = True
            _set_text(tf, slide_data.get("title", ""), 28, bold=True, color=_DARK)

            # Bullets
            bullets = slide_data.get("bullets", [])
            if bullets:
                bullet_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.6), Inches(5.5))
                tf2 = bullet_box.text_frame
                tf2.word_wrap = True
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        p = tf2.paragraphs[0]
                    else:
                        p = tf2.add_paragraph()
                    p.text = f"• {bullet}"
                    p.space_before = Pt(8)
                    for run in p.runs:
                        run.font.size = Pt(20)
                        run.font.color.rgb = _DARK

        # Speaker notes
        note = slide_data.get("speaker_note", "")
        if note:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = note

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


async def generate_presentation(documents: list[DocumentData]) -> bytes:
    slides_data = await generate_slides_data(documents)
    doc_title = documents[0].filename.rsplit(".", 1)[0] if documents else "Document"
    return build_pptx(slides_data, doc_title)
