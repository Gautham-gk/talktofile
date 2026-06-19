"""
Orchestrator Agent — manages the document processing pipeline state machine.
Pipeline: Extract → Lingua (translate) → Analyst (embed+index) → Ready
"""

import asyncio
import io
from enum import Enum
from typing import Callable, Awaitable

from openai import AsyncOpenAI

from agents.lingua_agent import detect_only
from agents.analyst_agent import analyse_one, generate_suggested_questions, generate_multi_doc_questions
from core.config import get_settings
from core.session_store import DocumentSession, DocumentData

TABULAR_EXTS = {"xlsx", "xls", "csv"}


class PipelineStage(str, Enum):
    EXTRACTING = "extracting"
    TRANSLATING = "translating"
    ANALYSING = "analysing"
    READY = "ready"
    ERROR = "error"


class NoReadableTextError(ValueError):
    """Raised when a document has no extractable text layer (e.g. a scanned PDF)."""


class MalformedFileError(ValueError):
    """Raised when a file is corrupt or cannot be parsed as its declared type."""


def _join_pages(pages: list[str]) -> str:
    """Join page texts with explicit page markers so the assistant can locate and
    return specific pages when a user asks (e.g. "give me pages 2-3")."""
    out = []
    for i, t in enumerate(pages, 1):
        if t and t.strip():
            out.append(f"[Page {i}]\n{t.strip()}")
    return "\n\n".join(out)


def _extract_pdf_text(content: bytes) -> str:
    """Extract text from a PDF, trying pdfplumber first then PyMuPDF.

    Different engines succeed on different PDFs (CID/embedded fonts common in
    European documents often defeat one but not the other), so we try both and
    keep whichever yields more text. Pages are tagged with [Page N] markers.
    """
    best_pages: list[str] = []

    # Engine 1: pdfplumber
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            best_pages = [page.extract_text() or "" for page in pdf.pages]
    except Exception:
        best_pages = []

    # Engine 2: PyMuPDF (fitz) — often recovers text pdfplumber misses.
    try:
        import fitz  # PyMuPDF
        with fitz.open(stream=content, filetype="pdf") as doc:
            alt_pages = [page.get_text("text") or "" for page in doc]
        if len("".join(alt_pages).strip()) > len("".join(best_pages).strip()):
            best_pages = alt_pages
    except Exception:
        pass

    best = _join_pages(best_pages)
    if not best.strip():
        raise NoReadableTextError(
            "This PDF has no embedded text layer — it looks like a scanned or "
            "image-only document. Please upload a text-based PDF (or run OCR on it first)."
        )
    return best


# Zip-based formats (docx/xlsx) can be "zip bombs" — tiny files that decompress
# to many GB. Reject anything whose declared uncompressed size is implausible.
_MAX_UNCOMPRESSED_BYTES = 200 * 1024 * 1024  # 200 MB


def _guard_zip_bomb(content: bytes) -> None:
    import zipfile
    try:
        with zipfile.ZipFile(io.BytesIO(content)) as zf:
            total = sum(info.file_size for info in zf.infolist())
    except zipfile.BadZipFile:
        raise MalformedFileError("File is not a valid Office document (bad archive).")
    if total > _MAX_UNCOMPRESSED_BYTES:
        raise MalformedFileError("File is too large when uncompressed and was rejected.")


def _extract_raw(filename: str, content: bytes) -> str:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext == "pdf":
        return _extract_pdf_text(content)

    elif ext == "docx":
        _guard_zip_bomb(content)
        from docx import Document
        doc = Document(io.BytesIO(content))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    elif ext in ("xlsx", "xls"):
        _guard_zip_bomb(content)
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        rows = []
        for sheet in wb.worksheets:
            rows.append(f"=== Sheet: {sheet.title} ===")
            for row in sheet.iter_rows(values_only=True):
                row_str = "\t".join(str(c) if c is not None else "" for c in row)
                if row_str.strip():
                    rows.append(row_str)
        return "\n".join(rows)

    elif ext == "pptx":
        _guard_zip_bomb(content)
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    texts.append(shape.text_frame.text.strip())
            if texts:
                slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides)

    elif ext in ("html", "htm"):
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(content.decode("utf-8", errors="replace"), "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        return soup.get_text(separator="\n").strip()

    elif ext == "json":
        import json as _json
        try:
            parsed = _json.loads(content.decode("utf-8", errors="replace"))
            return _json.dumps(parsed, indent=2, ensure_ascii=False)
        except _json.JSONDecodeError:
            raise MalformedFileError("File is not valid JSON.")

    else:
        # txt / csv / md plus source-code & other plain-text formats — read as text.
        return content.decode("utf-8", errors="replace")


def extract_text_from_file(filename: str, content: bytes) -> str:
    """Extract text, translating parser failures into clear, user-facing errors.

    Raises NoReadableTextError (empty / scanned) or MalformedFileError (corrupt)
    so the pipeline can stop the session with a meaningful message instead of crashing.
    """
    if not content:
        raise MalformedFileError(f"'{filename}' is empty — there's nothing to read.")

    try:
        text = _extract_raw(filename, content)
    except (NoReadableTextError, ValueError):
        raise
    except Exception as e:  # corrupt archive, bad zip, broken PDF structure, etc.
        raise MalformedFileError(
            f"'{filename}' appears to be corrupted or is not a valid file of its type "
            f"and couldn't be read. ({type(e).__name__})"
        )

    if not text.strip():
        raise NoReadableTextError(f"No readable text was found in '{filename}'.")
    return text


ProgressCallback = Callable[[PipelineStage, str], Awaitable[None]]


async def _analyse_into_document(filename: str, content: bytes, client: AsyncOpenAI) -> DocumentData:
    """Extract → detect language → chunk/embed/index/summarise one file."""
    raw_text = extract_text_from_file(filename, content)
    lang_code, _ = await detect_only(raw_text)
    chunks, embeddings, index, summary = await analyse_one(raw_text, client)
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    return DocumentData(
        filename=filename,
        original_language=lang_code,
        chunks=chunks,
        embeddings=embeddings,
        index=index,
        summary=summary,
        raw_text=raw_text,
        is_tabular=ext in TABULAR_EXTS,
    )


async def run_pipeline(
    session: DocumentSession,
    files: list[tuple[str, bytes]],
    on_progress: ProgressCallback,
) -> None:
    """Process one or more files into the session.

    A single malformed/unreadable file aborts the whole session with a clear message.
    """
    settings = get_settings()
    client = AsyncOpenAI(api_key=settings.openai_api_key)

    try:
        documents: list[DocumentData] = []
        total = len(files)

        for i, (filename, content) in enumerate(files, 1):
            label = f"({i}/{total}) " if total > 1 else ""
            await on_progress(PipelineStage.EXTRACTING, f"{label}Analysing {filename}...")
            doc = await _analyse_into_document(filename, content, client)
            documents.append(doc)

        session.documents = documents

        # Mode-aware suggested questions / actions.
        await on_progress(PipelineStage.ANALYSING, "Preparing suggestions...")
        if total == 1:
            session.suggested_questions = await generate_suggested_questions(
                documents[0].raw_text, client
            )
        else:
            from agents.analyst_agent import summary_to_text
            summaries = [(d.filename, summary_to_text(d.summary)) for d in documents]
            session.suggested_questions = await generate_multi_doc_questions(
                summaries, session.mode, client
            )

        session.ready = True

        # Note: the authoritative "ready" event (with summaries, languages, and
        # suggested questions) is sent by the document router after this returns.

    except Exception as e:
        session.ready = False
        await on_progress(PipelineStage.ERROR, str(e))
        raise
